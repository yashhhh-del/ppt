# ai_ppt_generator_ultimate.py
"""
AI PPT Generator — Ultimate (GPT-5 + Bing + Modern Gradient theme + extras)

Features:
- Uses OpenAI GPT-5 for slide outline, bullets, and speaker notes (via OPENAI_API_KEY)
- Uses Bing Web Search API (if BING_SUBSCRIPTION_KEY set) or SerpAPI or HTML-scrape fallback
- robots.txt respect (basic)
- newspaper3k + BeautifulSoup article extraction
- Caching to disk
- Unsplash image fallback
- Auto numeric detection & chart generation
- Modern Gradient theme for slides
- Speaker notes, slide timing, references
- Export PPTX and optional PDF conversion (if libreoffice available)
- Streamlit GUI + CLI mode
"""

import os
import re
import io
import json
import time
import hashlib
import logging
import requests
import subprocess
from urllib.parse import urlparse, quote_plus
from pathlib import Path
from collections import Counter
from typing import List, Dict

from bs4 import BeautifulSoup
from newspaper import Article
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

import matplotlib.pyplot as plt
from PIL import Image

from joblib import Memory
from cachetools import TTLCache
from dotenv import load_dotenv

# Streamlit UI
import streamlit as st

# Optional libs
try:
    import openai
except Exception:
    openai = None

nltk.download("punkt", quiet=True)

# -------------------------
# Config & logging
# -------------------------
load_dotenv()
logger = logging.getLogger("ai_ppt_ultimate")
logger.setLevel(logging.INFO)

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "").strip()
BING_KEY = os.getenv("BING_SUBSCRIPTION_KEY", "").strip()
BING_ENDPOINT = os.getenv("BING_ENDPOINT", "https://api.bing.microsoft.com/v7.0/search")
SERPAPI_KEY = os.getenv("SERPAPI_KEY", "").strip()
UNSPLASH_KEY = os.getenv("UNSPLASH_ACCESS_KEY", "").strip()

# caching
cache_dir = ".ai_ppt_cache"
memory = Memory(location=cache_dir, verbose=0)
ttl_cache = TTLCache(maxsize=1024, ttl=60*60*6)

HEADERS = {"User-Agent":"Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}

# -------------------------
# Utilities
# -------------------------
def safe_filename(s: str, ext="pptx"):
    s2 = re.sub(r'[^A-Za-z0-9_\- ]+', '', s).strip().replace(' ','_')[:120]
    return f"{s2}.{ext}"

def domain_from_url(u: str):
    try:
        return urlparse(u).netloc
    except:
        return ""

def check_robots_txt(url: str) -> bool:
    # Basic robots check - if robots.txt disallows "/" for user-agent: *
    try:
        domain = domain_from_url(url)
        if not domain:
            return True
        robots_url = f"https://{domain}/robots.txt"
        r = requests.get(robots_url, headers=HEADERS, timeout=6)
        if r.status_code != 200:
            return True
        txt = r.text.lower()
        if "user-agent: *" in txt:
            part = txt.split("user-agent: *",1)[1]
            disallows = re.findall(r'disallow:\s*(\S+)', part)
            for d in disallows:
                if d.strip() == "/":
                    return False
        return True
    except Exception as e:
        logger.warning("robots check error: %s", e)
        return True

# -------------------------
# Search layer (Bing -> SerpAPI -> HTML)
# -------------------------
def bing_search(query: str, count: int=5) -> List[str]:
    if not BING_KEY:
        return []
    try:
        headers = {"Ocp-Apim-Subscription-Key": BING_KEY}
        params = {"q": query, "count": count}
        r = requests.get(BING_ENDPOINT, headers=headers, params=params, timeout=10)
        r.raise_for_status()
        data = r.json()
        links = []
        for it in data.get("webPages", {}).get("value", [])[:count]:
            url = it.get("url")
            if url:
                links.append(url)
        return links
    except Exception as e:
        logger.warning("Bing search failed: %s", e)
        return []

def serpapi_search(query: str, num: int=5) -> List[str]:
    if not SERPAPI_KEY:
        return []
    try:
        params = {"engine":"google","q":query,"api_key":SERPAPI_KEY,"num":num}
        r = requests.get("https://serpapi.com/search.json", params=params, timeout=12)
        data = r.json()
        links = []
        for ritem in data.get("organic_results", [])[:num]:
            link = ritem.get("link')") or ritem.get("link") or ritem.get("url")
            if link:
                links.append(link)
        return links
    except Exception as e:
        logger.warning("SerpAPI failed: %s", e)
        return []

def html_bing_scrape(query: str, max_results: int=5) -> List[str]:
    q = quote_plus(query)
    url = f"https://www.bing.com/search?q={q}"
    try:
        r = requests.get(url, headers=HEADERS, timeout=10)
        soup = BeautifulSoup(r.text, "html.parser")
        links = []
        for a in soup.select("li.b_algo h2 a"):
            href = a.get('href')
            if href and href.startswith("http"):
                links.append(href)
            if len(links) >= max_results:
                break
        return links
    except Exception as e:
        logger.warning("HTML search scrape failed: %s", e)
        return []

def search_top_links(query: str, num: int=5) -> List[str]:
    # priority: Bing API -> SerpAPI -> HTML scrape
    links = []
    if BING_KEY:
        links = bing_search(query, count=num)
    if not links and SERPAPI_KEY:
        links = serpapi_search(query, num=num)
    if not links:
        links = html_bing_scrape(query, max_results=num)
    return links[:num]

# -------------------------
# Article extraction (cached)
# -------------------------
@memory.cache
def extract_article(url: str) -> Dict:
    try:
        if not check_robots_txt(url):
            return {"title": url, "text": "", "image": None, "url": url, "disallowed": True}
        art = Article(url)
        art.download()
        art.parse()
        return {"title": art.title or url, "text": art.text or "", "image": getattr(art,'top_image', None), "url": url, "disallowed": False}
    except Exception:
        try:
            r = requests.get(url, headers=HEADERS, timeout=10)
            soup = BeautifulSoup(r.text, "html.parser")
            ps = soup.find_all("p")
            text = "\n".join(p.get_text().strip() for p in ps)
            title = (soup.find("title").get_text() if soup.find("title") else url)
            og = soup.find('meta', property='og:image') or soup.find('meta', attrs={'name':'og:image'})
            top_image = og['content'] if og and og.get('content') else None
            return {"title": title, "text": text, "image": top_image, "url": url, "disallowed": False}
        except Exception as e:
            logger.warning("extract_article fallback failed: %s", e)
            return {"title": url, "text": "", "image": None, "url": url, "disallowed": False}

# -------------------------
# Summarization & slide outlines using OpenAI GPT-5
# -------------------------
def openai_init():
    if not OPENAI_API_KEY:
        return False
    if not openai:
        logger.warning("openai lib not installed")
        return False
    openai.api_key = OPENAI_API_KEY
    return True

def openai_chat_completion(messages: List[Dict], model="gpt-5", max_tokens=1200, temperature=0.15):
    if not openai_init():
        raise RuntimeError("OpenAI not configured")
    # ChatCompletion API usage; adapt if your account uses different call
    try:
        resp = openai.ChatCompletion.create(model=model, messages=messages, max_tokens=max_tokens, temperature=temperature)
        return resp['choices'][0]['message']['content'].strip()
    except Exception as e:
        logger.error("OpenAI request error: %s", e)
        raise

def generate_slide_structure(topic: str, aggregated_text: str, max_slides: int=10):
    # ask GPT-5 for JSON outline: slide_no, title, bullets[], notes
    prompt = f"""
You are an expert presentation designer. Create a slide deck for the topic below.
Constraints:
- Total slides (excluding title and references) <= {max_slides}
- Each slide: title, 3-6 concise bullet points (each <= 80 chars), and 2-4 short speaker notes sentences.
- Style: professional, modern, action-oriented.
- Use the aggregated research text to extract facts; where numbers are used, prefer accuracy. Provide a "Key Takeaways" slide near the end.
Output strictly as valid JSON array like:
[{{"slide_no":1,"title":"...","bullets":["..."],"notes":"..."}}, ...]
Topic: {topic}

Aggregated Research Text (reference):
\"\"\"{aggregated_text[:5000]}\"\"\"
"""
    try:
        out = openai_chat_completion([{"role":"user","content":prompt}], model="gpt-5", max_tokens=1400)
        # try parse JSON out of response
        try:
            return json.loads(out)
        except json.JSONDecodeError:
            # extract first JSON-looking block
            m = re.search(r'(\[.*\])', out, flags=re.S)
            if m:
                return json.loads(m.group(1))
    except Exception as e:
        logger.warning("OpenAI generation failed: %s", e)
    # Fallback simple extractive outline
    return fallback_outline(topic, aggregated_text, max_slides)

def fallback_outline(topic: str, text: str, max_slides: int=8):
    slides = []
    slides.append({"slide_no":1, "title":topic, "bullets":[topic], "notes":"Auto-generated title slide."})
    paras = [p for p in re.split(r'\n{2,}', text) if len(p.strip())>80]
    idx = 2
    for p in paras[:max_slides-3]:
        sents = sent_tokenize(p)[:3]
        bullets = [s.strip()[:80] for s in sents]
        slides.append({"slide_no":idx, "title": bullets[0][:40], "bullets":bullets, "notes":"Discuss these points."})
        idx += 1
    slides.append({"slide_no":idx, "title":"Key Takeaways", "bullets":["Summarize key points."], "notes":"Wrap up."})
    return slides

# -------------------------
# Numeric detection + chart
# -------------------------
def extract_numbers(text: str):
    patterns = re.findall(r'[\₹\$\€]?\d{1,3}(?:[,\.\d]{0,})\s?(?:billion|million|bn|m|crore|%|percent|years|year)?', text, flags=re.I)
    counts = {}
    for p in patterns:
        k = p.strip()
        counts[k] = counts.get(k,0)+1
    return counts

def make_stats_chart(numbers: Dict):
    if not numbers:
        return None
    items = sorted(numbers.items(), key=lambda x: x[1], reverse=True)[:8]
    labels = [i[0] for i in items][::-1]
    values = [i[1] for i in items][::-1]
    fig, ax = plt.subplots(figsize=(8,3.5))
    ax.barh(range(len(labels)), values)
    ax.set_yticks(range(len(labels)))
    ax.set_yticklabels(labels)
    ax.set_xlabel("Mentions")
    ax.set_title("Detected numeric mentions")
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf

# -------------------------
# Image fetching (OG image or Unsplash)
# -------------------------
def fetch_image(url: str):
    try:
        r = requests.get(url, headers=HEADERS, timeout=10)
        if r.status_code == 200:
            return io.BytesIO(r.content)
    except Exception as e:
        logger.warning("fetch_image failed: %s", e)
    return None

def unsplash_image(query: str):
    if not UNSPLASH_KEY:
        return None
    try:
        headers = {"Accept-Version":"v1", "Authorization":f"Client-ID {UNSPLASH_KEY}"}
        params = {"query": query, "per_page":1}
        r = requests.get("https://api.unsplash.com/search/photos", headers=headers, params=params, timeout=8)
        if r.status_code == 200:
            data = r.json()
            if data.get("results"):
                url = data["results"][0]["urls"]["regular"]
                return fetch_image(url)
    except Exception as e:
        logger.warning("Unsplash failed: %s", e)
    return None

# -------------------------
# PPTX creation with Modern Gradient theme
# -------------------------
def apply_modern_gradient_theme(prs: Presentation):
    # python-pptx can't change master easily; we'll style text boxes per slide
    return

def add_title_slide(prs: Presentation, title: str, subtitle: str=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    try:
        if subtitle:
            slide.placeholders[1].text = subtitle
    except:
        pass
    return slide

def add_content_slide(prs: Presentation, title: str, bullets: List[str], notes: str=None, image_bytes: io.BytesIO=None):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if notes:
        try:
            ns = slide.notes_slide
            ns.notes_text_frame.text = notes
        except:
            pass
    if image_bytes:
        try:
            left = Inches(6.0)
            top = Inches(1.2)
            slide.shapes.add_picture(image_bytes, left, top, width=Inches(3.0))
        except Exception:
            pass
    # Style accents (title color gradient impression via solid color)
    try:
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(33, 37, 111)
    except:
        pass
    return slide

def add_image_slide(prs: Presentation, title: str, image_bytes: io.BytesIO, notes: str=None):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    try:
        slide.shapes.title.text = title
    except:
        pass
    try:
        left = Inches(0.5)
        top = Inches(1.0)
        slide.shapes.add_picture(image_bytes, left, top, width=Inches(9.0))
    except Exception:
        pass
    if notes:
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except:
            pass
    return slide

def create_presentation(topic: str, slides_struct: List[Dict], articles: List[Dict], stats_img_buf: io.BytesIO=None, out_filename: str=None):
    if not out_filename:
        out_filename = safe_filename(topic, ext="pptx")
    prs = Presentation()
    apply_modern_gradient_theme(prs)
    add_title_slide(prs, topic, subtitle="Auto-generated — AI PPT Generator (Ultimate)")
    # add slides
    for s in slides_struct:
        title = s.get("title","Slide")
        bullets = s.get("bullets", [])[:6]
        notes = s.get("notes","")
        # try to find relevant article image
        img_buf = None
        for a in articles:
            if a.get("title") and title.lower() in (a.get("title","").lower()):
                if a.get("image"):
                    img_buf = fetch_image(a.get("image"))
                    break
        if not img_buf:
            img_buf = unsplash_image(title)
        add_content_slide(prs, title, bullets, notes=notes, image_bytes=img_buf)
    if stats_img_buf:
        add_image_slide(prs, "Detected Figures & Numbers", stats_img_buf, notes="Auto-extracted numeric mentions")
    # references
    refs = [a['url'] for a in articles if a.get('url')]
    if refs:
        add_content_slide(prs, "References & Sources", refs[:10], notes="Sources used to create this presentation")
    prs.save(out_filename)
    return out_filename

# -------------------------
# PDF conversion (if libreoffice available)
# -------------------------
def convert_pptx_to_pdf(pptx_path: str):
    # uses soffice --headless --convert-to pdf <file> --outdir <dir>
    out_dir = os.path.abspath(os.path.dirname(pptx_path))
    try:
        subprocess.run(["soffice", "--headless", "--convert-to", "pdf", pptx_path, "--outdir", out_dir], check=True, timeout=60)
        pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
        if os.path.exists(pdf_path):
            return pdf_path
    except Exception as e:
        logger.warning("PDF conversion failed: %s", e)
    return None

# -------------------------
# High-level pipeline
# -------------------------
def run_pipeline(topic: str, num_sources: int=4, max_slides:int=10):
    # 1) gather links
    links = search_top_links(topic, num=num_sources)
    articles = []
    for l in links:
        art = extract_article(l)
        articles.append(art)
    aggregated_text = "\n\n".join([a.get('text','') for a in articles if a.get('text')])
    # 2) generate slides using GPT-5 (or fallback)
    slides_struct = []
    if OPENAI_API_KEY and openai:
        try:
            slides_struct = generate_slide_structure(topic, aggregated_text, max_slides=max_slides)
        except Exception as e:
            logger.warning("OpenAI slide gen failed: %s; using fallback", e)
            slides_struct = fallback_outline(topic, aggregated_text, max_slides)
    else:
        slides_struct = fallback_outline(topic, aggregated_text, max_slides)
    # 3) numeric chart
    numbers = extract_numbers(aggregated_text)
    stats_img = make_stats_chart(numbers) if numbers else None
    # 4) create PPTX
    out_file = create_presentation(topic, slides_struct, articles, stats_img_buf=stats_img, out_filename=None)
    # 5) convert to PDF if possible
    pdf_file = convert_pptx_to_pdf(out_file)
    return {"pptx": out_file, "pdf": pdf_file, "slides": slides_struct, "sources": links}

# -------------------------
# Streamlit UI
# -------------------------
def streamlit_app():
    st.set_page_config(page_title="AI PPT Generator — Ultimate", layout="wide")
    st.title("AI PPT Generator — Ultimate (GPT-5 + Bing + Modern Gradient)")
    st.markdown("Generate a high-quality, professional PPTX from a topic. Uses GPT-5 for slide phrasing and speaker notes. Provide API keys as environment variables.")
    with st.sidebar:
        st.markdown("**Config / Keys**")
        st.write("- OPENAI_API_KEY: required for GPT-5.")
        st.write("- BING_SUBSCRIPTION_KEY: recommended for fast search.")
        st.write("- UNSPLASH_ACCESS_KEY: optional for images.")
        st.markdown("---")
        num_sources = st.slider("Number of sources", min_value=2, max_value=8, value=4)
        max_slides = st.slider("Max slides (excluding title/references)", min_value=4, max_value=20, value=10)
        do_pdf = st.checkbox("Try convert PPTX → PDF (requires LibreOffice)", value=False)
    topic = st.text_input("Enter presentation topic", value="Artificial Intelligence in Healthcare (2025 Trends)")
    if st.button("Generate Ultimate PPT"):
        if not OPENAI_API_KEY:
            st.warning("OPENAI_API_KEY not set. Will use extractive fallback (less polish).")
        st.info("Running pipeline — network requests will be made. This may take a while.")
        try:
            res = run_pipeline(topic, num_sources=num_sources, max_slides=max_slides)
            st.success(f"Generated: {res['pptx']}")
            with open(res['pptx'], "rb") as f:
                st.download_button("Download PPTX", data=f, file_name=res['pptx'], mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            if res.get("pdf"):
                with open(res['pdf'], "rb") as f:
                    st.download_button("Download PDF", data=f, file_name=os.path.basename(res['pdf']), mime="application/pdf")
            st.markdown("### Slide Outline")
            for s in res['slides']:
                st.markdown(f"**{s.get('title')}**")
                for b in s.get('bullets', [])[:6]:
                    st.write("- " + b)
                st.write("*Speaker notes:* " + (s.get('notes') or ""))
            st.markdown("### Sources")
            for src in res['sources']:
                st.write("- " + src)
        except Exception as e:
            st.error(f"Failed: {e}")
            logger.exception(e)
    else:
        st.info("Enter a topic and click Generate Ultimate PPT.")

# -------------------------
# CLI fallback
# -------------------------
def cli_main():
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--topic", required=True)
    parser.add_argument("--sources", type=int, default=4)
    parser.add_argument("--slides", type=int, default=10)
    args = parser.parse_args()
    out = run_pipeline(args.topic, num_sources=args.sources, max_slides=args.slides)
    print("PPTX:", out["pptx"])
    if out.get("pdf"):
        print("PDF:", out["pdf"])
    print("Slides:")
    for s in out["slides"]:
        print("-", s.get("title"))

if __name__ == "__main__":
    # If launched under streamlit, streamlit will import, so run CLI only if direct
    if os.getenv("STREAMLIT_RUN") or os.getenv("RUN_STREAMLIT_APP"):
        streamlit_app()
    else:
        # default: try to run streamlit app (so use `streamlit run file.py`)
        try:
            streamlit_app()
        except Exception:
            cli_main()
